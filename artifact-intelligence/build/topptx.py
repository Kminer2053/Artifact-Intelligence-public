#!/usr/bin/env python3
"""완성된 발표 슬라이드 HTML → 편집 가능한 PPTX 전환.

genre=slides 전용이다. 텍스트·표·도형·테두리는 **네이티브(편집 가능) 요소**로 옮기고,
도식·픽토그램 아이콘·이미지는 라이브 크롬에서 그 자리를 그대로 오려 그림으로 넣는다.

읽는 눈은 build/화면읽기.py 하나다 — 크롬 열기·CDP 붙기·JS 평가·사각형 캡처를 그대로
빌려 쓴다(포트0 자동배정이라 측정 크롬끼리 안 겹치고, 크롬 행 회수도 그 코드에 들었다).
그래서 pdftoppm·PDF 왕복이 없다(핸드오프 PoC 의 '200dpi 페이지 렌더 후 크롭' 경로를
'그 네모만 직접 캡처'로 대체 — 함정 #1 크롬 행·#6 pdftoppm 경고가 통째로 사라진다).

이것은 **전환이지 생성이 아니다** — tohwpx 와 같은 사상이다. 완성된 규격(화면의 최종
배치·서체·색)을 PPTX 그릇에 그대로 옮긴다. 규칙으로 다시 만들지 않는다.

폰트 — Pretendard 는 PPTX 에 임베드할 수 없어 미설치 PC 에서 예측불가 대체가 됐다.
그래서 공공 Windows(발표 실물)에 100% 있는 **맑은 고딕**으로 박는다(사장님 판정
2026-08-16, 아래 FONT). eastAsia 축까지 같은 글꼴로 지정해 한글이 적용되게 한다.
"""
from __future__ import annotations

import base64
import io
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import urllib.request
from pathlib import Path

from PIL import Image

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))  # 화면읽기·크롬찾기 가 같은 build/
import 화면읽기
from 크롬찾기 import 찾기

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

FONT = "맑은 고딕"   # pptx 는 Windows PowerPoint 에서 열어 고치는 편집형이라 그 PC 에 100%
                    # 있는 공공 표준 글꼴로 박는다(사장님 판정 2026-08-16). Pretendard 는
                    # 임베드 불가라 미설치 PC 에서 예측불가 대체가 됐다. 화면(HTML)은 그대로
                    # Pretendard — pptx 만 맑은 고딕(획·비례가 가까워 인상 유지).
_캡처배율 = 3   # 크롭 스크린샷 해상도(CSS px × 3 ≈ 인쇄 해상도). clip.scale 로만 준다.

# ── 화면에서 뽑는 수확기 ────────────────────────────────────────────────────
# harvest.js v2 를 CDP 반환형으로 옮긴 것. 요소 좌표·서체·색에 더해 테두리·의사요소
# (::before/::after)·CSS content 텍스트까지 뽑는다 — 헤드 세로바(border-left)·구분선
# (border-bottom)·카드 상단바(border-top)·"출처: " 접두(content)가 전부 이 두 계열이라
# 사장님 실물 1차 확인('26-08-14)에서 이게 빠져 재작업했다. 크롭(도식·아이콘·이미지)은
# 자리(ax/ay 문서좌표)만 실어 보내고, 파이썬이 그 자리를 화면읽기._찍기 로 오려 담는다.
_수확코드 = r"""(function () {
  const INLINE = new Set(["B","STRONG","I","EM","SPAN","BR","SMALL","SUB","SUP"]);
  function borderSides(cs) {
    const out = {};
    for (const nm of ["Top","Right","Bottom","Left"]) {
      const w = parseFloat(cs["border"+nm+"Width"]) || 0;
      const st = cs["border"+nm+"Style"];
      if (w >= 0.25 && st !== "none" && st !== "hidden")
        out[nm[0].toLowerCase()] = { w: w, c: cs["border"+nm+"Color"] };
    }
    return Object.keys(out).length ? out : null;
  }
  function pseudoBox(el, which, pr) {
    const cs = getComputedStyle(el, which);
    if (!cs.content || cs.content === "none" || cs.content === "normal") return null;
    const m = cs.content.match(/^"(.*)"$/);
    if (m && m[1]) return null;              // 알맹이 있는 텍스트만 텍스트 쪽 — content:"" 는 장식
    if (cs.position !== "absolute") return null;
    let cb = el;
    while (cb && getComputedStyle(cb).position === "static") cb = cb.parentElement;
    const cbr = (cb || document.body).getBoundingClientRect();
    const W = parseFloat(cs.width) || 0, H = parseFloat(cs.height) || 0;
    if (W < 1 || H < 1) return null;
    let x = cbr.left, y = cbr.top;
    if (cs.left !== "auto") x = cbr.left + parseFloat(cs.left);
    else if (cs.right !== "auto") x = cbr.right - parseFloat(cs.right) - W;
    if (cs.top !== "auto") y = cbr.top + parseFloat(cs.top);
    else if (cs.bottom !== "auto") y = cbr.bottom - parseFloat(cs.bottom) - H;
    const br = cs.borderRadius || "0";
    return { kind: "box", x: x - pr.left, y: y - pr.top, w: W, h: H,
             bg: cs.backgroundColor, bgi: cs.backgroundImage === "none" ? "" : cs.backgroundImage,
             radius: parseFloat(br) || 0, radiusPct: /%/.test(br),
             opacity: parseFloat(cs.opacity) };
  }
  function pseudoText(el, which) {
    const c = getComputedStyle(el, which).content || "";
    const m = c.match(/^"(.*)"$/);
    return m ? m[1] : "";
  }
  const pages = Array.from(document.querySelectorAll(".sl-page"));
  const out = [];
  for (const pg of pages) {
    const pr = pg.getBoundingClientRect();
    const pcs = getComputedStyle(pg);
    const items = [];
    const captured = [];
    items.push({ kind: "page", w: pr.width, h: pr.height,
                 bg: pcs.backgroundColor, bgi: pcs.backgroundImage });
    for (const which of ["::before","::after"]) {
      const b = pseudoBox(pg, which, pr);
      if (b) items.push(b);
    }
    const all = pg.querySelectorAll("*");
    for (const el of all) {
      const tag = el.tagName.toUpperCase();
      if (tag === "SCRIPT" || tag === "STYLE") continue;
      if (el.closest("svg") && tag !== "SVG") continue;
      if (el.closest("table") && tag !== "TABLE") continue;
      const r = el.getBoundingClientRect();
      if (r.width < 1 || r.height < 1) continue;
      const cs = getComputedStyle(el);
      if (cs.display === "none" || cs.visibility === "hidden") continue;
      const rel = { x: r.left - pr.left, y: r.top - pr.top, w: r.width, h: r.height };

      if (tag === "TABLE") {
        const rows = Array.from(el.rows).map(tr => Array.from(tr.cells).map(td => {
          const tcs = getComputedStyle(td);
          const tr2 = td.getBoundingClientRect();
          return { t: td.textContent.replace(/\s+/g, " ").trim(),
                   w: tr2.width, h: tr2.height, th: td.tagName === "TH",
                   bg: tcs.backgroundColor, color: tcs.color,
                   size: parseFloat(tcs.fontSize), b: (+tcs.fontWeight) >= 600,
                   align: tcs.textAlign, bd: borderSides(tcs) };
        }));
        items.push(Object.assign({ kind: "table", rows }, rel));
        continue;
      }
      // 도식(svgfig 주입 SVG)·이미지 — 자리만 실어 보내고 파이썬이 그 네모를 직접 찍는다.
      // svg 는 늘 진한 내용(도식·픽토, 테마색)이 있어 백지로 잡히면 실패로 본다(svg 표식).
      // img(첨부)는 밝을 수 있어 백지 검사에서 뺀다.
      if (tag === "SVG") {
        items.push(Object.assign({ kind: "crop", svg: true, ax: r.left, ay: r.top }, rel));
        continue;
      }
      if (tag === "IMG") {
        items.push(Object.assign({ kind: "crop", ax: r.left, ay: r.top }, rel));
        continue;
      }

      const bBox = pseudoBox(el, "::before", pr);
      if (bBox) items.push(bBox);

      const hasBg = cs.backgroundColor && cs.backgroundColor !== "rgba(0, 0, 0, 0)"
                    && cs.backgroundColor !== "transparent";
      const hasGrad = cs.backgroundImage && cs.backgroundImage !== "none";
      if (hasBg || hasGrad)
        items.push(Object.assign({ kind: "box", bg: cs.backgroundColor,
          bgi: hasGrad ? cs.backgroundImage : "",
          radius: parseFloat(cs.borderRadius) || 0,
          radiusPct: /%/.test(cs.borderRadius || ""), opacity: 1 }, rel));

      const bd = borderSides(cs);
      if (bd) items.push(Object.assign({ kind: "borders", bd }, rel));

      const directText = Array.from(el.childNodes).some(
        n => n.nodeType === 3 && n.textContent.trim());
      const onlyInlineKids = Array.from(el.children).every(c => INLINE.has(c.tagName.toUpperCase()));
      // <br>은 공백 없이 지우면 줄이 붙는다("제목<br>부제" → "제목부제"). 형제 수확기
      // (화면읽기.py)와 같은 규정 — br 있으면 클론에서 공백으로 바꿔 뽑는다(원본 DOM 은
      // 안 건드린다, 다른 요소 측정 오염 방지).
      let txt;
      if (el.querySelector("br")) {
        const c = el.cloneNode(true);
        c.querySelectorAll("br").forEach(b => b.replaceWith(" "));
        txt = c.textContent.replace(/\s+/g, " ").trim();
      } else {
        txt = el.textContent.replace(/\s+/g, " ").trim();
      }
      if (txt && directText && onlyInlineKids && !captured.some(a => a.contains(el))) {
        captured.push(el);
        let pre = pseudoText(el, "::before");
        let post = pseudoText(el, "::after");
        const par = el.parentElement;
        if (par && par !== pg) {
          if (el === par.firstElementChild) pre = pseudoText(par, "::before") + pre;
          if (el === par.lastElementChild) post = post + pseudoText(par, "::after");
        }
        txt = pre + txt + post;
        items.push(Object.assign({ kind: "text", t: txt,
          size: parseFloat(cs.fontSize), weight: +cs.fontWeight,
          color: cs.color, align: cs.textAlign,
          lh: parseFloat(cs.lineHeight) || 0, ff: cs.fontFamily }, rel));
      }

      const aBox = pseudoBox(el, "::after", pr);
      if (aBox) items.push(aBox);
    }
    out.push(items);
  }
  return JSON.stringify(out);
})()"""


# ── 색·칠·테두리 헬퍼 (build_pptx.py 포팅) ─────────────────────────────────
def _rgb(s):
    m = re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?", s or "")
    if not m:
        return None
    # 완전 투명(alpha 0)은 **색이 아니다** — None 을 줘 대체(흰색·무채움)로 가게 한다.
    # 안 그러면 "rgba(0, 0, 0, 0)"(배경 없는 표 셀)이 앞 세 숫자로 검정이 돼 셀을
    # 불투명 검정으로 칠하고 어두운 글씨가 사라진다(적대검토 HIGH, '26-08-14 실측).
    if m.group(4) is not None and float(m.group(4)) == 0:
        return None
    return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))


def _rgba_alpha(s):
    m = re.match(r"rgba?\(\d+,\s*\d+,\s*\d+(?:,\s*([\d.]+))?", s or "")
    return float(m.group(1)) if m and m.group(1) is not None else 1.0


def _grad_colors(bgi):
    # 크롬 computed style 은 color-mix() 결과를 color(srgb r g b) 0~1 실수로 직렬화한다
    cols = []
    for m in re.finditer(
            r"rgba?\((\d+),\s*(\d+),\s*(\d+)|color\(srgb\s+([\d.]+)\s+([\d.]+)\s+([\d.]+)",
            bgi or ""):
        if m.group(1) is not None:
            cols.append(RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3))))
        else:
            cols.append(RGBColor(*(round(float(m.group(i)) * 255) for i in (4, 5, 6))))
    return cols


def _no_line(shape):
    shape.line.fill.background()


def _set_grad(fill, colors, angle_deg=45):
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = colors[0]
    stops[-1].color.rgb = colors[-1]
    try:
        fill.gradient_angle = angle_deg
    except Exception:
        pass


def _solid_alpha(sh, color, alpha):
    """반투명 단색 칠 — python-pptx 공개 API 엔 알파가 없어 spPr 에 직접 쓴다.
    스키마 순서(prstGeom→fill→ln) 때문에 기존 fill·ln 을 걷어내고 fill 을 먼저 넣는다."""
    spPr = sh._element.spPr
    for t in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:ln"):
        for e in spPr.findall(qn(t)):
            spPr.remove(e)
    fill = spPr.makeelement(qn("a:solidFill"), {})
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": str(color)})
    a = spPr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
    clr.append(a)
    fill.append(clr)
    spPr.append(fill)
    ln = spPr.makeelement(qn("a:ln"), {})
    ln.append(spPr.makeelement(qn("a:noFill"), {}))
    spPr.append(ln)   # ln 을 안 쓰면 테마 윤곽선이 상속된다


def _cell_borders(cell, bd, k):
    """표 셀 테두리 — tcPr 의 lnL·lnR·lnT·lnB (스키마상 fill 앞이라 맨 앞에 끼운다)"""
    tcPr = cell._tc.get_or_add_tcPr()
    order = [("l", "a:lnL"), ("r", "a:lnR"), ("t", "a:lnT"), ("b", "a:lnB")]
    made = []
    for key, tag in order:
        side = bd.get(key)
        if not side:
            continue
        ln = tcPr.makeelement(qn(tag), {"w": str(int(side["w"] * k * 12700)), "cap": "flat"})
        f = tcPr.makeelement(qn("a:solidFill"), {})
        c = tcPr.makeelement(qn("a:srgbClr"), {"val": str(_rgb(side["c"]) or RGBColor(0xD9, 0xD9, 0xD9))})
        f.append(c)
        ln.append(f)
        made.append(ln)
    for ln in reversed(made):
        tcPr.insert(0, ln)


_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}


def _거의흴가(png_bytes, 문턱=0.995):
    """캡처가 사실상 백지(내용 없음)인가 — 48×48 로 줄여 흰 픽셀 비율로 본다. svg 도식·
    픽토는 늘 진한 테마색 내용이 있어 99.5% 이상 흰색이면 렌더 실패로 본다(적대검토 #7 —
    좌표가 밀려 빈 자리를 찍으면 유효한 흰 PNG 가 나와 png=None 고발을 우회했다)."""
    im = Image.open(io.BytesIO(png_bytes)).convert("L").resize((48, 48))
    d = im.tobytes()   # "L" 모드 = 픽셀당 1바이트(밝기). getdata() 는 Pillow 14 에서 사라진다.
    return sum(1 for v in d if v > 245) >= len(d) * 문턱


# ── 수확 (라이브 크롬) ──────────────────────────────────────────────────────
def 수확(html경로) -> tuple[list | None, str | None]:
    """HTML 을 머리 없는 크롬으로 열어 페이지별 요소를 뽑고, 크롭 자리는 그 자리 그대로
    찍어 png(base64)를 실어 돌려준다. 실패는 (None, 사유)."""
    크롬 = 찾기()   # 못 찾아도 죽지 않는다(WP-S8) — 서버가 이 요청 하나로 안 내려가게
    if not 크롬:
        return None, "크롬을 찾지 못해 PPTX 를 뽑을 수 없습니다"
    with tempfile.TemporaryDirectory() as tmp:
        p = subprocess.Popen(
            [크롬, "--headless", "--disable-gpu", "--remote-debugging-port=0",
             f"--user-data-dir={tmp}/u", "--no-first-run", "--no-default-browser-check",
             Path(html경로).resolve().as_uri()],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        try:
            붙을곳 = None
            for _ in range(160):
                try:
                    항 = int((Path(tmp) / "u" / "DevToolsActivePort")
                             .read_text().splitlines()[0])
                    j = json.load(urllib.request.urlopen(f"http://127.0.0.1:{항}/json", timeout=1))
                    쓸것 = [t for t in j if t.get("type") == "page" and t.get("webSocketDebuggerUrl")]
                    if 쓸것:
                        붙을곳 = 쓸것[0]["webSocketDebuggerUrl"]
                        break
                except Exception:
                    pass
                time.sleep(0.25)
            if not 붙을곳:
                return None, "크롬 디버깅에 못 붙었습니다"
            # 인쇄 매체를 **수확 전에 따로** 앉힌다. 이걸 수확 eval 의 앞선것으로 같이 보내면
            # 미디어 재적용이 그 eval 안에서 아직 정착 안 된 배치를 재게 해 좌표가 밀린다(도식
            # y 3345→3194 로 밀려 5쪽 도식이 백지로 캡처됨, '26-08-14 사장님 실물). 화면에는
            # 발표 조작 UI(present.js: 앞뒤·PDF·목차)가 같이 뜨는데 그건 문서가 아니다(@media
            # print 에서 display:none). PDF 도 같은 매체라 결과가 맞는다.
            화면읽기._평가(붙을곳, "1", [("Emulation.setEmulatedMedia", {"media": "print"})])
            time.sleep(1.3)   # 글꼴·자간·svgfig 도식·미디어 전환이 다 앉을 때까지
            pages = json.loads(화면읽기._평가(붙을곳, _수확코드, []))   # 미디어 재적용 없이 잰다
            # 도식·아이콘·이미지 — 그 네모만 그대로 찍는다(배율 3, 화면읽기._찍기 와 같은 눈).
            for items in pages:
                for it in items:
                    if it.get("kind") == "crop":
                        # **직접 clip 캡처 — setDeviceMetricsOverride 를 안 쓴다.** 화면읽기._찍기 는
                        # 캡처마다 그걸 걸어 뷰포트를 바꿔 문서를 리플로우시켰다: 아래쪽 장일수록
                        # 좌표가 밀려(도식 y 3194→3345) 수확 때 잰 자리엔 백지만 잡혔다('26-08-14
                        # 사장님 실물에서 5쪽 도식이 빈 것으로 드러남). 레이아웃을 안 건드리면 수확
                        # 좌표가 그대로 유효하다 — clip.scale 로만 해상도를 올린다.
                        받 = 화면읽기._명령(붙을곳, "Page.captureScreenshot", {
                            "format": "png", "captureBeyondViewport": True,
                            "clip": {"x": it["ax"], "y": it["ay"], "width": it["w"],
                                     "height": it["h"], "scale": _캡처배율}})
                        it["png"] = (받 or {}).get("data")
            return pages, None
        except (Exception, SystemExit) as e:
            # 화면읽기 는 웹소켓 끊김·수확 JS 예외를 **SystemExit**(BaseException)로 던진다.
            # 그게 위로 새면 동기 호출자(api.부르기·serve._post·verify_all)의 `except Exception`
            # 을 뚫고 ok:False 계약이 깨진다 — 스레드 서버는 워커만 죽어 응답이 안 나가고
            # verify_all 은 통째로 중단된다(적대검토 MEDIUM). 여기서 붙잡아 실패로 돌린다.
            return None, f"수확 중 크롬이 끊겼습니다: {str(e)[:80]}"
        finally:
            p.terminate()
            try:
                p.wait(timeout=5)
            except Exception:
                p.kill()


# ── 조립 (수확 → PPTX) ─────────────────────────────────────────────────────
def _조립(pages, tmpdir):
    """수확한 페이지들을 960×540pt 슬라이드로 조립한다. 크롭 png 가 비면(도식을 못
    옮겼으면) 조용히 넘기지 않고 목록에 담아 돌려준다 — 부른 쪽이 실패로 고발한다."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)    # 960pt
    prs.slide_height = Emu(6858000)    # 540pt
    blank = prs.slide_layouts[6]
    빠진크롭 = []

    for pageno, items in enumerate(pages, 1):
        page = next(it for it in items if it["kind"] == "page")
        k = 960.0 / page["w"]          # px → pt
        s = prs.slides.add_slide(blank)

        # 페이지 배경 — 슬라이드 배경 속성은 뷰어가 그라데이션을 무시하는 일이 있어
        # (실측: 검정) 전면 사각형 도형으로 깐다. 도형 gradFill 은 표준대로 먹는다.
        cols = _grad_colors(page.get("bgi", ""))
        bgc = _rgb(page.get("bg", ""))
        if len(cols) >= 2 or bgc is not None:
            bgsh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                      prs.slide_width, prs.slide_height)
            if len(cols) >= 2:
                _set_grad(bgsh.fill, cols)
            else:
                bgsh.fill.solid()
                bgsh.fill.fore_color.rgb = bgc
            _no_line(bgsh)
            bgsh.shadow.inherit = False

        for idx, it in enumerate(items):
            kind = it["kind"]
            if kind == "page":
                continue
            x, y, w, h = (Pt(it["x"] * k), Pt(it["y"] * k),
                          Pt(max(it["w"] * k, 1)), Pt(max(it["h"] * k, 1)))

            if kind == "box":
                # 반경을 '짧은 변 대비 비율'로 환산한다 — %는 /100, px는 짧은 변으로 나눈다.
                # 옛 코드는 "20%"의 20 을 px 처럼 min(w,h)/2 와 견줘 작은 상자를 타원으로
                # 오판했다(적대검토 LOW). 비율이 0.5 이상이면 타원, 그 아래 양수면 둥근 사각형.
                rad = it["radius"]
                미니 = min(it["w"], it["h"]) or 1
                비율 = (rad / 100.0) if it.get("radiusPct") else (rad / 미니)
                shape_kind = (MSO_SHAPE.OVAL if 비율 >= 0.5
                              else MSO_SHAPE.ROUNDED_RECTANGLE if rad > 0
                              else MSO_SHAPE.RECTANGLE)
                sh = s.shapes.add_shape(shape_kind, x, y, w, h)
                if shape_kind == MSO_SHAPE.ROUNDED_RECTANGLE:
                    try:
                        sh.adjustments[0] = min(비율, 0.5)   # 실제 곡률(기본 16.7% 고정 탈피)
                    except Exception:
                        pass
                cols = _grad_colors(it.get("bgi", ""))
                alpha = _rgba_alpha(it.get("bg", "")) * it.get("opacity", 1)
                if len(cols) >= 2:
                    _set_grad(sh.fill, cols)
                    _no_line(sh)
                elif alpha < 0.999:
                    c = _rgb(it.get("bg", "")) or RGBColor(0xFF, 0xFF, 0xFF)
                    _solid_alpha(sh, c, alpha)     # fill·ln 을 스스로 정리한다
                else:
                    c = _rgb(it.get("bg", "")) or RGBColor(0xFF, 0xFF, 0xFF)
                    sh.fill.solid()
                    sh.fill.fore_color.rgb = c
                    _no_line(sh)
                sh.shadow.inherit = False

            elif kind == "borders":
                bd = it["bd"]
                px, py, pw, ph = it["x"], it["y"], it["w"], it["h"]
                for key, bx, by, bw, bh in (
                        ("t", px, py, pw, bd.get("t", {}).get("w", 0)),
                        ("b", px, py + ph - bd.get("b", {}).get("w", 0), pw, bd.get("b", {}).get("w", 0)),
                        ("l", px, py, bd.get("l", {}).get("w", 0), ph),
                        ("r", px + pw - bd.get("r", {}).get("w", 0), py, bd.get("r", {}).get("w", 0), ph)):
                    if key not in bd:
                        continue
                    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Pt(bx * k), Pt(by * k),
                                            Pt(max(bw * k, 0.5)), Pt(max(bh * k, 0.5)))
                    sh.fill.solid()
                    sh.fill.fore_color.rgb = _rgb(bd[key]["c"]) or RGBColor(0xD9, 0xD9, 0xD9)
                    _no_line(sh)
                    sh.shadow.inherit = False

            elif kind == "crop":
                b64 = it.get("png")
                png = base64.b64decode(b64) if b64 else None
                # svg 크롭(도식·픽토)이 백지로 잡히면 렌더 실패다 — 조용히 빈 그림을 심지
                # 않고 고발한다(적대검토 #7: '26-08-14 사장님 실물에서 5쪽 도식이 백지로
                # 나갔다). img(첨부)는 밝을 수 있어 백지 검사에서 뺀다.
                if not png or (it.get("svg") and _거의흴가(png)):
                    빠진크롭.append(f"{pageno}쪽 도식/그림")
                    continue
                fn = os.path.join(tmpdir, f"crop-{pageno}-{idx}.png")
                with open(fn, "wb") as f:
                    f.write(png)
                s.shapes.add_picture(fn, x, y, w, h)

            elif kind == "text":
                tb = s.shapes.add_textbox(x, y, w, h)
                tf = tb.text_frame
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                # HTML 에서 한 줄이던 텍스트는 PPT 에서도 꺾지 않는다 — 대체 글꼴이 조금
                # 넓어도 상자 밖으로 흐르게 두는 쪽이 강제 줄바꿈보다 낫다(표지 실측).
                줄수 = round(it["h"] / it["lh"]) if it.get("lh") else 1
                tf.word_wrap = 줄수 > 1
                p = tf.paragraphs[0]
                p.alignment = _ALIGN.get(it.get("align", "left"), PP_ALIGN.LEFT)
                if it.get("lh"):
                    p.line_spacing = Pt(it["lh"] * k)
                r = p.add_run()
                r.text = it["t"]
                r.font.name = FONT
                r.font.size = Pt(round(it["size"] * k * 2) / 2)
                r.font.bold = it.get("weight", 400) >= 600
                c = _rgb(it.get("color", ""))
                if c is not None:
                    r.font.color.rgb = c
                rPr = r._r.get_or_add_rPr()   # 한글 글꼴은 eastAsia 축도 지정해야 적용된다
                rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": FONT}))

            elif kind == "table":
                rows = it["rows"]
                ncol = max(len(r_) for r_ in rows)
                tbl = s.shapes.add_table(len(rows), ncol, x, y, w, h).table
                for j, cw in enumerate(rows[0]):
                    tbl.columns[j].width = Pt(cw["w"] * k)
                for i2, row in enumerate(rows):
                    tbl.rows[i2].height = Pt(row[0]["h"] * k)
                    for j, cell in enumerate(row):
                        tc = tbl.cell(i2, j)
                        tc.margin_left = tc.margin_right = Pt(6)
                        tc.margin_top = tc.margin_bottom = Pt(3)
                        if cell.get("bd"):
                            _cell_borders(tc, cell["bd"], k)
                        bg = _rgb(cell.get("bg", ""))
                        tc.fill.solid()
                        tc.fill.fore_color.rgb = bg if bg is not None else RGBColor(0xFF, 0xFF, 0xFF)
                        p = tc.text_frame.paragraphs[0]
                        p.alignment = _ALIGN.get(cell.get("align", "left"), PP_ALIGN.LEFT)
                        r = p.add_run()
                        r.text = cell["t"]
                        r.font.name = FONT
                        r.font.size = Pt(round(cell["size"] * k * 2) / 2)
                        r.font.bold = cell.get("b", False)
                        c = _rgb(cell.get("color", ""))
                        if c is not None:
                            r.font.color.rgb = c
                        rPr = r._r.get_or_add_rPr()
                        rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": FONT}))

    return prs, 빠진크롭


def _되읽기(경로):
    """편집 가능성 구조 되읽기 — 텍스트가 텍스트인지, 표가 네이티브 표인지 센다."""
    p2 = Presentation(경로)
    n = {"슬라이드": len(p2.slides), "텍스트": 0, "표": 0, "그림": 0}
    for sl in p2.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                n["텍스트"] += 1
            if sh.has_table:
                n["표"] += 1
            if sh.shape_type == 13:   # PICTURE
                n["그림"] += 1
    return n


def 만들기(html경로, 낼경로) -> tuple[bool, object]:
    """슬라이드 HTML → PPTX. tohwpx.만들기 와 같은 (ok, 말) 계약.
    성공: (True, 되읽기요약dict). 실패: (False, 사유문자열)."""
    if not os.path.exists(html경로):
        return False, "HTML 이 없습니다 — 조립을 먼저 돌리세요"
    pages, 사유 = 수확(html경로)
    if 사유:
        return False, 사유
    if not pages:
        return False, "슬라이드를 하나도 읽지 못했습니다"
    tmpdir = tempfile.mkdtemp(prefix="pptxcrop_")
    try:
        prs, 빠진크롭 = _조립(pages, tmpdir)
        if 빠진크롭:
            # 조용히 빈 자리로 두지 않는다 — 도식을 못 찍었으면 실패로 알린다(화면읽기 철학).
            return False, f"도식·그림 {len(빠진크롭)}건을 옮기지 못했습니다: {', '.join(빠진크롭[:3])}"
        prs.save(낼경로)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    return True, _되읽기(낼경로)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("사용법: python topptx.py <슬라이드.html> <낼.pptx>", file=sys.stderr)
        sys.exit(2)
    ok, 말 = 만들기(sys.argv[1], sys.argv[2])
    print(("OK " if ok else "실패 ") + json.dumps(말, ensure_ascii=False))
    sys.exit(0 if ok else 1)
